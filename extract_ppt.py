from pptx import Presentation
import json

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.replace("\n", " ").strip())
        slides_content.append({"slide_number": i + 1, "content": slide_text})
    return slides_content

if __name__ == "__main__":
    pptx_path = r"e:\Durgesh work\lesca-tech\lesca-tech-ppt.pptx"
    content = extract_text_from_pptx(pptx_path)
    with open("extracted_content.json", "w", encoding="utf-8") as f:
        json.dump(content, f, indent=4)
    print("Content extracted successfully to extracted_content.json")
